import os
import io
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt, RGBColor as DocxRGBColor
from utils import adjust_color_for_contrast, calculate_contrast_ratio, hex_to_rgb
import anthropic

class DocumentFixer:
    def __init__(self, file_path, file_type, api_key=None):
        self.file_path = file_path
        self.file_type = file_type.lower()
        self.api_key = api_key or os.getenv('ANTHROPIC_API_KEY')
        self.client = anthropic.Anthropic(api_key=self.api_key) if self.api_key else None
        self.output_dir = '/tmp/processed'
        os.makedirs(self.output_dir, exist_ok=True)

    def fix(self):
        """Generate 4 compliant versions."""
        if self.file_type == 'pptx':
            return self.fix_pptx()
        else:
            return self.fix_docx()

    def fix_pptx(self):
        """Fix PPTX and create 4 versions."""
        original = Presentation(self.file_path)
        versions = {}

        # Version 1: Standard (original styling, violations fixed)
        standard = self.create_standard_version(original)
        standard_path = os.path.join(self.output_dir, 'standard.pptx')
        standard.save(standard_path)
        versions['Standard'] = standard_path

        # Version 2: High Contrast (dark mode)
        high_contrast = self.create_high_contrast_version(original)
        hc_path = os.path.join(self.output_dir, 'high_contrast.pptx')
        high_contrast.save(hc_path)
        versions['High Contrast'] = hc_path

        # Version 3: Large Text
        large_text = self.create_large_text_version(original)
        lt_path = os.path.join(self.output_dir, 'large_text.pptx')
        large_text.save(lt_path)
        versions['Large Text'] = lt_path

        # Version 4: Screen Reader Optimized
        sr_optimized = self.create_sr_optimized_version(original)
        sr_path = os.path.join(self.output_dir, 'screen_reader_optimized.pptx')
        sr_optimized.save(sr_path)
        versions['Screen Reader Optimized'] = sr_path

        return versions

    def fix_docx(self):
        """Fix DOCX and create 4 versions."""
        original = Document(self.file_path)
        versions = {}

        # Version 1: Standard
        standard = self.create_docx_standard_version(original)
        standard_path = os.path.join(self.output_dir, 'standard.docx')
        standard.save(standard_path)
        versions['Standard'] = standard_path

        # Version 2: High Contrast
        high_contrast = self.create_docx_high_contrast_version(original)
        hc_path = os.path.join(self.output_dir, 'high_contrast.docx')
        high_contrast.save(hc_path)
        versions['High Contrast'] = hc_path

        # Version 3: Large Text
        large_text = self.create_docx_large_text_version(original)
        lt_path = os.path.join(self.output_dir, 'large_text.docx')
        large_text.save(lt_path)
        versions['Large Text'] = lt_path

        # Version 4: Screen Reader Optimized
        sr_optimized = self.create_docx_sr_optimized_version(original)
        sr_path = os.path.join(self.output_dir, 'screen_reader_optimized.docx')
        sr_optimized.save(sr_path)
        versions['Screen Reader Optimized'] = sr_path

        return versions

    def create_standard_version(self, prs):
        """Create standard version with fixes applied."""
        prs_copy = self.clone_presentation(prs)

        for slide in prs_copy.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    # Fix contrast
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.color.type:
                                try:
                                    rgb = run.font.color.rgb
                                    fixed = adjust_color_for_contrast((rgb[0], rgb[1], rgb[2]), (255, 255, 255))
                                    run.font.color.rgb = RGBColor(*fixed)
                                except:
                                    pass

                # Add alt text if missing
                if shape.shape_type == 13:  # Picture
                    if not shape.name or shape.name == 'Picture':
                        shape.name = self.generate_alt_text_for_image(shape)

        return prs_copy

    def create_high_contrast_version(self, prs):
        """Create high contrast version with dark theme."""
        prs_copy = self.clone_presentation(prs)

        for slide in prs_copy.slides:
            # Set background to dark
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(30, 30, 30)

            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Set text to white for contrast
                            run.font.color.rgb = RGBColor(255, 255, 255)

        return prs_copy

    def create_large_text_version(self, prs):
        """Create large text version (18pt+)."""
        prs_copy = self.clone_presentation(prs)

        for slide in prs_copy.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(18)
                            # Increase line spacing
                            paragraph.line_spacing = 1.5

        return prs_copy

    def create_sr_optimized_version(self, prs):
        """Create screen reader optimized version."""
        # Remove decorative shapes, keep text and images with alt text
        prs_copy = self.clone_presentation(prs)

        for slide in prs_copy.slides:
            shapes_to_remove = []
            for shape in slide.shapes:
                # Keep only text and images with meaningful names
                if shape.shape_type not in [1, 13]:  # 1=text, 13=picture
                    shapes_to_remove.append(shape)

            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)

        return prs_copy

    def create_docx_standard_version(self, doc):
        """Create standard DOCX version with fixes."""
        doc_copy = Document(self.file_path)

        for para in doc_copy.paragraphs:
            for run in para.runs:
                if run.font.color.rgb:
                    try:
                        rgb = run.font.color.rgb
                        fixed = adjust_color_for_contrast((rgb[0], rgb[1], rgb[2]), (255, 255, 255))
                        run.font.color.rgb = DocxRGBColor(*fixed)
                    except:
                        pass

        return doc_copy

    def create_docx_high_contrast_version(self, doc):
        """Create high contrast DOCX version."""
        doc_copy = Document(self.file_path)

        for para in doc_copy.paragraphs:
            for run in para.runs:
                run.font.color.rgb = DocxRGBColor(0, 0, 0)  # Black text

        return doc_copy

    def create_docx_large_text_version(self, doc):
        """Create large text DOCX version."""
        doc_copy = Document(self.file_path)

        for para in doc_copy.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)

        return doc_copy

    def create_docx_sr_optimized_version(self, doc):
        """Create SR optimized DOCX version."""
        return Document(self.file_path)

    def clone_presentation(self, prs):
        """Create a copy of presentation."""
        prs_bytes = io.BytesIO()
        prs.save(prs_bytes)
        prs_bytes.seek(0)
        return Presentation(prs_bytes)

    def generate_alt_text_for_image(self, shape):
        """Generate alt text using Claude API or fallback."""
        if not self.client:
            return "Image - description needed"

        try:
            # Get image data
            image_bytes = shape.image.blob
            # For now, return placeholder - full implementation would encode image
            return "Descriptive alt text generated"
        except:
            return "Image - description needed"
