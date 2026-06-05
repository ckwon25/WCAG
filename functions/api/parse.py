from pptx import Presentation
from docx import Document
from PIL import Image
import io
import os

class DocumentParser:
    def __init__(self, file_path):
        self.file_path = file_path
        self.file_type = file_path.lower().split('.')[-1]
        self.content = {
            'type': self.file_type,
            'slides': [],
            'images': [],
            'colors': [],
            'headings': []
        }

    def parse(self):
        if self.file_type == 'pptx':
            return self.parse_pptx()
        elif self.file_type in ['docx', 'doc']:
            return self.parse_docx()
        else:
            raise ValueError(f"Unsupported file type: {self.file_type}")

    def parse_pptx(self):
        prs = Presentation(self.file_path)

        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                'index': slide_idx,
                'text_elements': [],
                'images': [],
                'shapes': []
            }

            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, 'text') and shape.text.strip():
                    text_info = {
                        'text': shape.text,
                        'color': self.extract_color(shape),
                        'font_size': self.extract_font_size(shape)
                    }
                    slide_data['text_elements'].append(text_info)

                # Extract images
                if shape.shape_type == 13:  # Picture
                    try:
                        image_data = self.extract_image_from_shape(shape)
                        if image_data:
                            slide_data['images'].append(image_data)
                            self.content['images'].append(image_data)
                    except:
                        pass

            self.content['slides'].append(slide_data)

        return self.content

    def parse_docx(self):
        doc = Document(self.file_path)

        for para in doc.paragraphs:
            if para.text.strip():
                level = para.style.name
                heading_level = 0

                # Detect heading level
                if level.startswith('Heading'):
                    heading_level = int(level.split()[-1])

                element = {
                    'text': para.text,
                    'type': 'heading' if heading_level > 0 else 'paragraph',
                    'level': heading_level,
                    'color': self.extract_color_from_paragraph(para),
                    'font_size': self.extract_font_size_from_paragraph(para)
                }

                if heading_level > 0:
                    self.content['headings'].append(element)

                self.content['slides'].append(element)

        # Extract images from document
        for rel in doc.part.rels.values():
            if 'image' in rel.target_ref:
                try:
                    image_data = self.extract_image_from_rel(rel)
                    if image_data:
                        self.content['images'].append(image_data)
                except:
                    pass

        return self.content

    def extract_color(self, shape):
        try:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.color.type:
                            rgb = run.font.color.rgb
                            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        except:
            pass
        return None

    def extract_color_from_paragraph(self, para):
        try:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        except:
            pass
        return None

    def extract_font_size(self, shape):
        try:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            return run.font.size.pt
        except:
            pass
        return 12

    def extract_font_size_from_paragraph(self, para):
        try:
            for run in para.runs:
                if run.font.size:
                    return run.font.size.pt
        except:
            pass
        return 12

    def extract_image_from_shape(self, shape):
        try:
            image = shape.image
            return {
                'data': image.blob,
                'filename': image.filename,
                'has_alt_text': bool(shape.name)
            }
        except:
            return None

    def extract_image_from_rel(self, rel):
        try:
            image_data = rel.target_part.blob
            return {
                'data': image_data,
                'filename': rel.target_ref
            }
        except:
            return None
